from pptx import Presentation

prs = Presentation()

def add_slide(title, content, note):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = content
    slide.notes_slide.notes_text_frame.text = note

slides = [
("Software Process Structure and Importance",
 "Software Engineering Lecture\nReal-World Software Development Approach",
 "Assalam-o-Alaikum students. Aaj hum Software Process Structure aur us ki importance ko samjhenge."),

("Introduction",
 "Software is used in every industry\nApps, websites, systems\nStructured process required",
 "Software har jaga use ho raha hai aur process zaroori hai."),

("What is Software Process?",
 "Step-by-step development method\nOrganized workflow",
 "Software process step by step method hai."),

("Importance of Software Process",
 "Better quality\nLess errors\nTime saving\nCost saving",
 "Process se quality improve hoti hai."),

("Software Process Structure",
 "Requirement\nDesign\nImplementation\nTesting\nDeployment\nMaintenance",
 "Ye 6 main steps hotay hain."),

("Requirement Analysis",
 "Collect user needs\nDefine system goals",
 "Requirements mein user needs samajhtay hain."),

("System Design",
 "Architecture design\nDatabase + UI planning",
 "Design mein system blueprint banta hai."),

("Implementation",
 "Coding starts\nConvert design into software",
 "Implementation mein coding hoti hai."),

("Testing",
 "Find bugs\nImprove quality",
 "Testing mein errors check hotay hain."),

("Deployment",
 "Software release\nLive system",
 "Deployment mein software live hota hai."),

("Maintenance",
 "Bug fixing\nUpdates\nImprovements",
 "Maintenance continuous hoti hai."),

("Thank You",
 "Questions?",
 "Thank you students.")
]

for s in slides:
    add_slide(*s)

prs.save("Software_Process_Presentation.pptx")
print("PPT Created Successfully!")