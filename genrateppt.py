from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()

# ---------- STYLE FUNCTION ----------
def set_style(shape, size=32, bold=True, color=(0, 51, 102)):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(*color)

# ---------- ADD SLIDE FUNCTION ----------
def add_slide(title, content, note):
    slide_layout = prs.slide_layouts[1]  # Title + Content
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]

    title_shape.text = title
    content_shape.text = content

    # Notes (speaker script)
    slide.notes_slide.notes_text_frame.text = note

    return slide

# ---------- TITLE SLIDE ----------
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Software Process Structure & Importance"
slide.placeholders[1].text = "Modern Lecture Presentation"

# ---------- SLIDES ----------
add_slide(
"Software Process Overview",
"• Requirement Analysis\n• Design\n• Implementation\n• Testing\n• Deployment\n• Maintenance",
"Software process ke 6 main steps hotay hain jo real-world development mein use hotay hain."
)

add_slide(
"Requirement Analysis",
"• Understand user needs\n• Collect requirements\n• Define system goals",
"Is phase mein hum user se requirements collect karte hain."
)

add_slide(
"Design Phase",
"• System architecture\n• Database design\n• UI planning",
"Design phase mein system ka blueprint banaya jata hai."
)

add_slide(
"Implementation",
"• Coding starts\n• Convert design into software",
"Coding is phase mein hoti hai."
)

add_slide(
"Testing Phase",
"• Find bugs\n• Fix errors\n• Improve quality",
"Testing mein software ki errors check ki jati hain."
)

add_slide(
"Deployment & Maintenance",
"• Software release\n• Updates\n• Bug fixing",
"Software users ke liye live kiya jata hai aur baad mein updates aate hain."
)

# ---------- SAVE ----------
prs.save("Modern_Software_Process_Presentation.pptx")
print("PPT Created Successfully!")